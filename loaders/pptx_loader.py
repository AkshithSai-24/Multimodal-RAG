"""
PPTX Loader — python-pptx for text / notes / tables / embedded images.
Each slide is also rendered as a SLIDE modality chunk (text summary).
"""

from __future__ import annotations

import base64
import io
import zipfile
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Inches

from loaders.base_loader import BaseLoader
from models.document import DocumentChunk, Modality, SourceType


class PPTXLoader(BaseLoader):
    """Load text, tables, and images from a .pptx file."""

    async def load(self, source: str, **kwargs) -> List[DocumentChunk]:
        path = Path(source)
        chunks: List[DocumentChunk] = []
        source_name = path.name
        prs = Presentation(str(path))

        # ── Slide-level extraction ────────────────────────────────────────────
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts: List[str] = []
            notes_text = ""

            # Notes
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                notes_text = tf.text.strip() if tf else ""

            for shape in slide.shapes:
                # ── Text frames ───────────────────────────────────────────────
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_texts.append(text)

                # ── Tables ────────────────────────────────────────────────────
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    md = self._rows_to_markdown(rows)
                    if md:
                        chunks.append(
                            DocumentChunk(
                                content=md,
                                modality=Modality.TABLE,
                                source_type=SourceType.PPTX,
                                source_id=str(path),
                                source_name=source_name,
                                page_number=slide_num,
                            )
                        )

            # ── Slide summary chunk ───────────────────────────────────────────
            full_text = "\n".join(slide_texts)
            if notes_text:
                full_text += f"\n\n[Speaker Notes]: {notes_text}"
            if full_text.strip():
                chunks.append(
                    DocumentChunk(
                        content=full_text,
                        modality=Modality.SLIDE,
                        source_type=SourceType.PPTX,
                        source_id=str(path),
                        source_name=source_name,
                        page_number=slide_num,
                        metadata={"has_notes": bool(notes_text)},
                    )
                )

        # ── Embedded images (via zip) ─────────────────────────────────────────
        with zipfile.ZipFile(str(path), "r") as zf:
            media_files = [n for n in zf.namelist() if n.startswith("ppt/media/")]
            for img_idx, media_path in enumerate(media_files):
                img_bytes = zf.read(media_path)
                if len(img_bytes) < 4096:
                    continue
                ext = media_path.rsplit(".", 1)[-1].lower()
                mime = f"image/{ext}" if ext in {"jpg", "jpeg", "png", "gif", "webp"} else "image/jpeg"
                b64 = base64.b64encode(img_bytes).decode()
                chunks.append(
                    DocumentChunk(
                        content=f"Slide image {img_idx + 1} from {source_name}",
                        modality=Modality.IMAGE,
                        source_type=SourceType.PPTX,
                        source_id=str(path),
                        source_name=source_name,
                        image_base64=b64,
                        image_mime_type=mime,
                        metadata={"image_index": img_idx},
                    )
                )

        return chunks

    @staticmethod
    def _rows_to_markdown(rows: List[List[str]]) -> str:
        if not rows:
            return ""
        header = "| " + " | ".join(rows[0]) + " |"
        sep = "| " + " | ".join(["---"] * len(rows[0])) + " |"
        body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
        return "\n".join(filter(None, [header, sep, body]))
